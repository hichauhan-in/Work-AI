"""Multi-format file loaders.

Each loader converts a file into a list of :class:`Section` objects (one per page,
slide, sheet, etc.) carrying provenance metadata. Heavy parsing libraries are imported
lazily inside each loader so this module imports cleanly without them and so missing an
optional parser only affects that one format.

Supported today:
  - Plain text / markdown / csv / log : .txt .md .markdown .csv .log
  - PDF (incl. OneNote-exported PDF)  : .pdf   (renders sparse pages for OCR)
  - Word / PowerPoint / Excel        : .docx .pptx .xlsx
  - HTML / RTF                       : .html .htm .rtf
  - Images (OCR + vision)            : .png .jpg .jpeg .bmp .tif .tiff .gif .webp
"""
from __future__ import annotations

import io
from pathlib import Path

from ..logging_setup import get_logger
from ..schema import Section

log = get_logger("loaders")

TEXT_EXT = {".txt", ".md", ".markdown", ".csv", ".log"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif", ".webp"}
DOC_EXT = {".pdf", ".docx", ".pptx", ".xlsx", ".html", ".htm", ".rtf"}
SUPPORTED_EXT = TEXT_EXT | IMAGE_EXT | DOC_EXT


def is_supported(path: str | Path) -> bool:
    return Path(path).suffix.lower() in SUPPORTED_EXT


def _base_meta(path: Path) -> dict:
    ext = path.suffix.lower()
    return {"source": str(path), "filename": path.name, "filetype": ext.lstrip(".")}


def load_file(path: str | Path, ocr=None) -> list[Section]:
    """Dispatch a file to the appropriate loader. Returns [] for unsupported types."""
    path = Path(path)
    ext = path.suffix.lower()
    meta = _base_meta(path)

    if ext in TEXT_EXT:
        return _load_text(path, meta)
    if ext == ".pdf":
        return _load_pdf(path, meta, ocr)
    if ext == ".docx":
        return _load_docx(path, meta, ocr)
    if ext == ".pptx":
        return _load_pptx(path, meta, ocr)
    if ext == ".xlsx":
        return _load_xlsx(path, meta)
    if ext in {".html", ".htm"}:
        return _load_html(path, meta)
    if ext == ".rtf":
        return _load_rtf(path, meta)
    if ext in IMAGE_EXT:
        return _load_image(path, meta, ocr)

    log.debug("Unsupported file skipped: %s", path)
    return []


# --------------------------------------------------------------------------- #
# Individual loaders
# --------------------------------------------------------------------------- #
def _load_text(path: Path, meta: dict) -> list[Section]:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return [Section(text, meta)] if text else []


def _load_pdf(path: Path, meta: dict, ocr) -> list[Section]:
    import fitz  # PyMuPDF, lazy

    sections: list[Section] = []
    doc = fitz.open(path)
    try:
        for page_index in range(doc.page_count):
            page = doc.load_page(page_index)
            text = page.get_text("text").strip()

            # Screenshot-heavy / OneNote pages often have little selectable text.
            # Render the whole page and OCR it so image content becomes searchable.
            if ocr and ocr.enabled and len(text.split()) < 10:
                ocr_text = _ocr_pdf_page(page, ocr)
                if ocr_text:
                    text = (text + "\n" + ocr_text).strip() if text else ocr_text

            if text:
                sections.append(
                    Section(text, {**meta, "page": page_index + 1, "pages": doc.page_count})
                )
    finally:
        doc.close()
    return sections


def _ocr_pdf_page(page, ocr) -> str:
    try:
        from PIL import Image  # lazy

        pix = page.get_pixmap(dpi=200)
        image = Image.open(io.BytesIO(pix.tobytes("png")))
        return ocr.pil_to_text(image)
    except Exception as exc:  # pragma: no cover
        log.warning("PDF page OCR failed: %s", exc)
        return ""


def _load_docx(path: Path, meta: dict, ocr) -> list[Section]:
    import docx  # python-docx, lazy

    document = docx.Document(str(path))
    parts: list[str] = [p.text for p in document.paragraphs if p.text and p.text.strip()]

    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    sections: list[Section] = []
    body = "\n".join(parts).strip()
    if body:
        sections.append(Section(body, meta))

    # OCR embedded images so screenshots pasted into Word are searchable.
    if ocr and ocr.enabled:
        image_text = _ocr_embedded_images(
            (rel.target_part.blob for rel in document.part.rels.values()
             if "image" in rel.reltype),
            ocr,
        )
        if image_text:
            sections.append(Section(image_text, {**meta, "content": "embedded_images"}))
    return sections


def _load_pptx(path: Path, meta: dict, ocr) -> list[Section]:
    from pptx import Presentation  # python-pptx, lazy

    presentation = Presentation(str(path))
    sections: list[Section] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        image_blobs = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if ocr and ocr.enabled and getattr(shape, "shape_type", None) == 13:  # PICTURE
                try:
                    image_blobs.append(shape.image.blob)
                except Exception:
                    pass
        if ocr and ocr.enabled and image_blobs:
            image_text = _ocr_embedded_images(iter(image_blobs), ocr)
            if image_text:
                texts.append(image_text)
        body = "\n".join(texts).strip()
        if body:
            sections.append(Section(body, {**meta, "slide": slide_index}))
    return sections


def _load_xlsx(path: Path, meta: dict) -> list[Section]:
    import openpyxl  # lazy

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    sections: list[Section] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            body = "\n".join(rows).strip()
            if body:
                sections.append(Section(body, {**meta, "sheet": sheet.title}))
    finally:
        workbook.close()
    return sections


def _load_html(path: Path, meta: dict) -> list[Section]:
    from bs4 import BeautifulSoup  # lazy

    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator="\n").strip()
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return [Section(text, meta)] if text else []


def _load_rtf(path: Path, meta: dict) -> list[Section]:
    from striprtf.striprtf import rtf_to_text  # lazy

    raw = path.read_text(encoding="utf-8", errors="replace")
    text = rtf_to_text(raw).strip()
    return [Section(text, meta)] if text else []


def _load_image(path: Path, meta: dict, ocr) -> list[Section]:
    if not (ocr and ocr.enabled):
        log.info("OCR disabled; skipping image %s", path)
        return []
    text = ocr.file_to_text(path)
    return [Section(text, {**meta, "content": "image_ocr"})] if text.strip() else []


def _ocr_embedded_images(blobs, ocr) -> str:
    """OCR an iterable of raw image bytes; returns concatenated text."""
    from PIL import Image  # lazy

    collected: list[str] = []
    for blob in blobs:
        try:
            with Image.open(io.BytesIO(blob)) as image:
                piece = ocr.pil_to_text(image)
            if piece:
                collected.append(piece)
        except Exception as exc:  # pragma: no cover
            log.debug("Embedded image OCR skipped: %s", exc)
    return "\n".join(collected).strip()
